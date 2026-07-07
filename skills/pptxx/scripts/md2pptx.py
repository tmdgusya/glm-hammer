"""md2pptx — runtime-only pptxx converter (M5, hardened from the M1 spike).

Hardens `planning/spikes/pptx/md2pptx-spike.py` into a deck-directory CLI:

    python md2pptx.py <deck-dir> [<tokens.json>]

- Parses `<deck-dir>/slides.md` (M4-era format-spec grammar).
- Resolves image / shots references ONLY inside `<deck-dir>` (incl. `<deck-dir>/shots/`);
  absolute, `file://` / `http(s)://`, `..`-escape, and non-image extensions are REFUSED
  with a nonzero exit (path allowlist — §0 "md2pptx path allowlist = deck directory").
- Diagram export policy: `<deck-dir>/shots/slide-<NN>.png` present -> exactly one picture
  shape; absent -> title only + description demoted into (non-empty) speaker notes, 0 pictures.
- Tokens are best-effort: the optional `<tokens.json>` arg, else `<deck-dir>/tokens.json`,
  else built-in defaults — a missing / unparseable tokens file never crashes (exit 0).
- Output: `<deck-dir>/<basename>.pptx` (overwrite allowed). Prints a slide-count summary.

RUNTIME ONLY: no hook or test imports / requires / spawns this file. `node tests/gates.test.js`
stays pure-Node and never invokes it; it only asserts this file exists + the zero-dep boundary.
"""
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from PIL import Image

SLIDE_W = Emu(12192000)  # 16:9
SLIDE_H = Emu(6858000)
MARGIN = Emu(838200)     # ~0.875in

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}

DEFAULT_TOKENS = {
    "surface": RGBColor.from_string("FFFFFF"),
    "text": RGBColor.from_string("111111"),
    "accent": RGBColor.from_string("B3541E"),
    "font": "Calibri",
    "size_title": Pt(30),
    "size_body": Pt(15),
}


def die(msg, code=2):
    sys.stderr.write("md2pptx: " + msg + "\n")
    sys.exit(code)


def refuse(msg):
    # Path-restriction rejection (C-3): nonzero exit + clear message.
    sys.stderr.write("md2pptx: REFUSED: " + msg + "\n")
    sys.exit(3)


def load_tokens(deck_dir, tokens_arg):
    """Best-effort token load — any failure falls back to DEFAULT_TOKENS (exit 0)."""
    path = tokens_arg or os.path.join(deck_dir, "tokens.json")
    try:
        with open(path, encoding="utf-8") as f:
            doc = json.load(f)
        hexv = lambda p: RGBColor.from_string(p["$value"].lstrip("#"))
        px = lambda p: Pt(float(p["$value"].rstrip("px")) * 0.75)  # 1px = 0.75pt
        return {
            "surface": hexv(doc["color"]["surface"]["default"]),
            "text": hexv(doc["color"]["text"]["default"]),
            "accent": hexv(doc["color"]["accent"]["default"]),
            "font": doc["typography"]["font"]["display"]["$value"][0],
            "size_title": px(doc["typography"]["size"]["title"]),
            "size_body": px(doc["typography"]["size"]["body"]),
        }
    except Exception:
        return dict(DEFAULT_TOKENS)


def resolve_image_ref(deck_dir, ref):
    """Validate an image reference against the deck-directory allowlist (C-3).

    Rejects (nonzero exit): absolute / drive / `file://` / `http(s)://` refs, `..`
    parent-escapes, non-image extensions, and any path normalising outside the deck.
    Returns an absolute path inside the deck subtree on success.
    """
    raw = (ref or "").strip()
    norm = raw.replace("\\", "/")
    if not raw:
        refuse("empty image reference")
    if norm.startswith("/"):
        refuse("absolute-path image ref rejected: " + ref)
    if re.match(r"^[A-Za-z]:", norm):
        refuse("absolute (drive-letter) image ref rejected: " + ref)
    if "://" in norm or norm.lower().startswith("file:"):
        refuse("URL / scheme image ref rejected: " + ref)
    if ".." in norm.split("/"):
        refuse("parent-escape ('..') image ref rejected: " + ref)
    ext = os.path.splitext(norm)[1].lower()
    if ext not in IMAGE_EXTS:
        refuse("non-image type image ref rejected: " + ref)
    deck_real = os.path.realpath(deck_dir)
    target = os.path.realpath(os.path.join(deck_dir, norm.replace("/", os.sep)))
    try:
        inside = os.path.commonpath([deck_real, target]) == deck_real
    except ValueError:
        inside = False
    if not inside:
        refuse("image ref escapes deck directory: " + ref)
    return target


def parse_slides(md_text):
    """Split on standalone `---` lines (fence-protected); extract each slide's parts."""
    chunks, cur, in_fence = [], [], False
    for line in md_text.splitlines():
        if line.strip().startswith("```"):
            in_fence = not in_fence
        if line == "---" and not in_fence:
            chunks.append(cur)
            cur = []
        else:
            cur.append(line)
    chunks.append(cur)

    slides = []
    for chunk in chunks:
        s = {"title": None, "body": [], "image": None, "attribution": None,
             "diagram": None, "notes": None}
        i = 0
        while i < len(chunk):
            line = chunk[i]
            if line.startswith("# ") and s["title"] is None:
                s["title"] = line[2:].strip()
            elif line == "```diagram":
                j = i + 1
                block = []
                while j < len(chunk) and chunk[j] != "```":
                    block.append(chunk[j])
                    j += 1
                m = re.match(r"slide:\s*(\d+)", block[0]) if block else None
                s["diagram"] = {
                    "slide": m.group(1) if m else None,
                    "text": "\n".join(block[1:] if m else block).strip(),
                }
                i = j
            elif line.startswith("> notes:"):
                s["notes"] = line[len("> notes:"):].strip()
            elif line.startswith("> attribution:"):
                s["attribution"] = line[len("> attribution:"):].strip()
            elif line.startswith("!["):
                m = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", line)
                if m:
                    s["image"] = m.group(2)
            elif line.strip():
                s["body"].append(line)
            i += 1
        slides.append(s)
    return slides


def add_text(slide, text, left, top, width, height, size, color, font, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    return box


def build(deck_dir, tk):
    with open(os.path.join(deck_dir, "slides.md"), encoding="utf-8") as f:
        slides_md = parse_slides(f.read())

    shots_dir = os.path.join(deck_dir, "shots")
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    for s in slides_md:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = tk["surface"]

        # Hardening guard (C-2): a title-less slide gets a placeholder, never a crash.
        title = s["title"] if s["title"] else "(untitled slide)"
        add_text(slide, title, MARGIN, Emu(600000), SLIDE_W - 2 * MARGIN,
                 Emu(900000), tk["size_title"], tk["text"], tk["font"], bold=True)

        notes = s["notes"]
        y = Emu(1700000)

        if s["diagram"]:
            nn = s["diagram"]["slide"]
            shot = os.path.join(shots_dir, "slide-{}.png".format(nn)) if nn else None
            if shot and os.path.exists(shot):
                # Default path: reuse the raster as exactly one picture shape (C-4).
                with Image.open(shot) as im:
                    w, h = im.size
                disp_h = Emu(4200000)
                disp_w = Emu(int(disp_h * w / h))
                slide.shapes.add_picture(shot, Emu(int((SLIDE_W - disp_w) / 2)), y,
                                         disp_w, disp_h)
            else:
                # Fallback (C-4): title only; demote the description into speaker notes,
                # guaranteed non-empty even when the description is blank.
                label = nn if nn else "?"
                demoted = "[diagram fallback — slide {}] ".format(label) \
                    + s["diagram"]["text"].replace("\n", " ")
                notes = (notes + " " + demoted) if notes else demoted
        elif s["image"]:
            img_path = resolve_image_ref(deck_dir, s["image"])  # C-3 allowlist
            if os.path.exists(img_path):
                pic_side = Emu(3600000)
                slide.shapes.add_picture(img_path, MARGIN, y, pic_side, pic_side)
                if s["attribution"]:
                    add_text(slide, "attribution: " + s["attribution"], MARGIN,
                             Emu(int(y + pic_side + 200000)), SLIDE_W - 2 * MARGIN,
                             Emu(500000), tk["size_body"], tk["accent"], tk["font"])
                for k, line in enumerate(s["body"]):
                    add_text(slide, line, Emu(int(MARGIN + pic_side + 400000)),
                             Emu(int(y + k * 500000)), Emu(5500000), Emu(500000),
                             tk["size_body"], tk["text"], tk["font"])
            else:
                # Valid but missing local image: best-effort, render body text only.
                for k, line in enumerate(s["body"]):
                    add_text(slide, line, MARGIN, Emu(int(y + k * 550000)),
                             SLIDE_W - 2 * MARGIN, Emu(550000),
                             tk["size_body"], tk["text"], tk["font"])
        else:
            for k, line in enumerate(s["body"]):
                add_text(slide, line, MARGIN, Emu(int(y + k * 550000)),
                         SLIDE_W - 2 * MARGIN, Emu(550000),
                         tk["size_body"], tk["text"], tk["font"])

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    return prs, len(slides_md)


def main():
    if len(sys.argv) < 2:
        die("usage: python md2pptx.py <deck-dir> [<tokens.json>]")
    deck_dir = sys.argv[1]
    tokens_arg = sys.argv[2] if len(sys.argv) > 2 else None
    if not os.path.isdir(deck_dir):
        die("deck directory not found: " + deck_dir)
    if not os.path.isfile(os.path.join(deck_dir, "slides.md")):
        die("slides.md not found in deck directory: " + deck_dir)

    tk = load_tokens(deck_dir, tokens_arg)
    prs, n = build(deck_dir, tk)

    basename = os.path.basename(os.path.normpath(deck_dir))
    out_path = os.path.join(deck_dir, basename + ".pptx")
    prs.save(out_path)
    print("wrote {} ({} slides)".format(out_path, n))


if __name__ == "__main__":
    main()
